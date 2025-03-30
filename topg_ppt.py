import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_topg_presentation(output_path='TOPG_Bot_Presentation.pptx'):
    """Create a PowerPoint presentation for TOPG Discord Bot"""
    
    # Create presentation with widescreen layout
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    colors = {
        'discord_blue': RGBColor(88, 101, 242),  # Discord blurple
        'background': RGBColor(45, 45, 48),      # Dark gray
        'text_primary': RGBColor(255, 255, 255),  # White
        'text_secondary': RGBColor(153, 170, 181),  # Light blue-gray
        'highlight': RGBColor(255, 200, 87),     # Yellow highlight
        'accent': RGBColor(114, 137, 218)        # Light discord blue
    }
    
    # Style settings for consistent formatting
    title_font_size = Pt(44)
    subtitle_font_size = Pt(32)
    heading_font_size = Pt(28)
    subheading_font_size = Pt(24)
    body_font_size = Pt(18)
    
    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Apply background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['discord_blue']
    
    # Add title and subtitle
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "TOPG Discord Bot"
    subtitle.text = "Crypto Tracking & Analysis Made Simple"
    
    # Format title text
    title_text_frame = title.text_frame
    title_paragraph = title_text_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_run = title_paragraph.runs[0]
    title_run.font.size = title_font_size
    title_run.font.color.rgb = colors['text_primary']
    title_run.font.bold = True
    
    # Format subtitle text
    subtitle_text_frame = subtitle.text_frame
    subtitle_paragraph = subtitle_text_frame.paragraphs[0]
    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    subtitle_run = subtitle_paragraph.runs[0]
    subtitle_run.font.size = subtitle_font_size
    subtitle_run.font.color.rgb = colors['text_primary']
    
    # Function to create a standard slide with title and content
    def create_content_slide(title_text, subtitle_text=None, content_items=None):
        slide_layout = prs.slide_layouts[1]  # Content slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Apply background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors['background']
        
        # Add title
        title = slide.shapes.title
        title.text = title_text
        title_text_frame = title.text_frame
        title_paragraph = title_text_frame.paragraphs[0]
        title_run = title_paragraph.runs[0]
        title_run.font.size = heading_font_size
        title_run.font.color.rgb = colors['discord_blue']
        title_run.font.bold = True
        
        # Add content placeholder
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        
        # Add subtitle if provided
        if subtitle_text:
            p = text_frame.paragraphs[0]
            p.text = subtitle_text
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.size = subheading_font_size
            run.font.color.rgb = colors['text_secondary']
            text_frame.add_paragraph()  # Add empty paragraph for spacing
        
        # Add content items if provided
        if content_items:
            # Check if first item is a section header
            current_paragraph = text_frame.paragraphs[0] if not subtitle_text else text_frame.add_paragraph()
            
            for item in content_items:
                # Check if item is a section header (starts with §)
                if item.startswith('§'):
                    section_p = text_frame.add_paragraph()
                    section_p.text = item[1:]  # Remove the § marker
                    section_p.alignment = PP_ALIGN.LEFT
                    section_run = section_p.runs[0]
                    section_run.font.size = subheading_font_size
                    section_run.font.color.rgb = colors['highlight']
                    section_run.font.bold = True
                    current_paragraph = text_frame.add_paragraph()
                # Check if item is a bullet point (starts with •)
                elif item.startswith('•'):
                    bullet_p = text_frame.add_paragraph()
                    bullet_p.text = item[1:].strip()  # Remove the • marker and whitespace
                    bullet_p.level = 1  # Set bullet level
                    run = bullet_p.runs[0]
                    run.font.size = body_font_size
                    run.font.color.rgb = colors['text_primary']
                # Regular paragraph
                else:
                    current_paragraph.text = item
                    run = current_paragraph.runs[0]
                    run.font.size = body_font_size
                    run.font.color.rgb = colors['text_primary']
                    current_paragraph = text_frame.add_paragraph()
        
        return slide
    
    # Introduction Slide
    intro_content = [
        "§Comprehensive Blockchain Coverage:",
        "• Track tokens across Solana, Ethereum, BNB, and Base networks with unified interface",
        "• Automated real-time information retrieval from multiple data sources",
        "• Intelligent message parsing for token addresses and ticker symbols",
        "§Powerful Features:",
        "• Highly configurable server-specific settings with granular controls",
        "• Performance-optimized asynchronous architecture for reliability",
        "• Database integration for persistent configuration and historical data",
        "• Advanced monitoring and self-maintenance capabilities"
    ]
    create_content_slide("Introduction", "Advanced Crypto Tracking for Discord Communities", intro_content)
    
    # Auto-Messaging System Slide
    auto_messaging_content = [
        "§Intelligent Message Parsing:",
        "• Automatically identifies and validates addresses across multiple blockchain formats",
        "• Detects token symbols with $ prefix and disambiguates when necessary",
        "• Processes multiple tokens mentioned in a single message with separate responses",
        "§Configurable Response Behavior:",
        "• Server-wide monitoring or channel-specific operation based on admin preference",
        "• Include or exclude specific channels from auto-messaging",
        "• Intelligent caching of frequently requested token data",
        "§Real-Time Data Integration:",
        "• Direct blockchain access through Solana RPC connections for on-chain data",
        "• Integration with multiple data providers for comprehensive market information",
        "• Fallback systems and retry logic for uninterrupted operation"
    ]
    create_content_slide("Auto-Messaging System", "Instant Token Information Without Commands", auto_messaging_content)
    
    # Token Information Display Slide
    token_info_content = [
        "§Financial Metrics:",
        "• Current USD value with precision formatting based on token value",
        "• FDV calculation and current market capitalization",
        "• Real-time liquidity pool depth measurement",
        "• 5-minute volume and price change tracking with visual indicators",
        "• All-time high market cap tracking with percentage from peak",
        "§Transaction Intelligence:",
        "• Recent transaction counts with buy/sell breakdown",
        "• Buy vs. sell pressure indicators for trend analysis",
        "• Detection and highlighting of significant token movements",
        "§Resource Integration:",
        "• Automatic discovery and linking of official token websites",
        "• Twitter and Telegram account linking for community access",
        "• Direct links to trading platforms (Axiom, Photon, Neo BullX)",
        "• DEX chart integration for immediate price analysis"
    ]
    create_content_slide("Token Information Display", "Comprehensive Market Data at a Glance", token_info_content)
    
    # First-Caller Tracking Slide
    first_caller_content = [
        "§User Recognition:",
        "• Automatic tracking of which user first mentions each token",
        "• Recording of exact time for transparency and verification",
        "• Username shown with each token information display for credit",
        "§Performance Tracking:",
        "• Records market cap at moment of first mention for baseline",
        "• Real-time calculation of return on investment since first mention",
        "• Visual multiplier indicators for significant growth milestones",
        "• Retention of performance data across all supported tokens",
        "§Verification Features:",
        "• Verification of token listing status on decentralized exchanges",
        "• Paid vs. organic listing indicators for complete context",
        "• Relative and absolute time tracking since first mention",
        "§Community Impact:",
        "• Encourages users to discover and share promising tokens early",
        "• Creates friendly competition for finding valuable projects",
        "• Allows users to develop track records as successful spotters"
    ]
    create_content_slide("First-Caller Tracking System", "Recognizing Early Token Discovery", first_caller_content)
    
    # Server Settings Management Slide
    settings_content = [
        "§Configuration Commands:",
        "• `/settings mode` - Switch between server-wide and channel-specific operating modes",
        "• `/settings add-channel` - Add specific channels to the auto-messaging whitelist",
        "• `/settings remove-channel` - Remove channels from auto-messaging to limit bot activity",
        "• `/settings status` - View detailed information about current configuration and active channels",
        "§Permission Security:",
        "• Configuration commands restricted to server administrators only",
        "• Flexible permission system based on Discord role hierarchy",
        "• Confirmation processes for sensitive configuration modifications",
        "§Operational Flexibility:",
        "• Server-Wide Mode: Monitor all channels with optional exclusion of specific channels",
        "• Channel-Specific Mode: Explicitly include only designated channels for targeted operation",
        "• Sensible defaults for new installations with easy customization"
    ]
    create_content_slide("Server Settings Management", "Flexible Configuration for Any Server", settings_content)
    
    # GitHub Repository Analysis Slide
    github_content = [
        "§Technical Evaluation:",
        "• Detailed scoring of codebase quality (out of 25 points)",
        "• Evaluation of project structure and organization",
        "• Assessment of coding practices and technical implementation",
        "• Scoring of documentation completeness and clarity",
        "§Investment Intelligence:",
        "• Overall score calculation based on multiple factors",
        "• Analysis of project credibility and transparency",
        "• Specialized evaluation of AI/ML implementation quality",
        "• Clear guidance with calculated confidence rating",
        "§Security Assessment:",
        "• Detection of potential security issues in codebase",
        "• Analysis of third-party library usage and risks",
        "• Verification of secure coding patterns and practices",
        "§Repository Insights:",
        "• Primary language identification and technology stack",
        "• Verification and explanation of project licensing",
        "• Creation and update timestamps for activity assessment",
        "• Statistics on stars, forks, and watchers for popularity"
    ]
    create_content_slide("GitHub Repository Analysis", "AI-Powered Project Assessment", github_content)
    
    # Health Monitoring System Slide
    health_content = [
        "§Performance Metrics:",
        "• Continuous monitoring of system availability with precise duration",
        "• Counter for total messages processed with rate analysis",
        "• Average processing time calculation for performance tuning",
        "• Measurement of response times across different components",
        "• Comprehensive error tracking with categorization and trending",
        "§API Performance:",
        "• Individual response time tracking for each API endpoint",
        "• Success rate calculation for external service dependencies",
        "• Status monitoring of third-party services and APIs",
        "§Resource Management:",
        "• Real-time tracking of memory usage with threshold alerts",
        "• Connection pool monitoring for optimal resource allocation",
        "• CPU and network utilization tracking for performance optimization"
    ]
    create_content_slide("Health Monitoring System", "Real-Time Performance Insights", health_content)
    
    # Technical Excellence Slide
    technical_content = [
        "§Asynchronous Architecture:",
        "• Non-blocking processing model for optimal responsiveness",
        "• Asynchronous operations for network and database interactions",
        "• Simultaneous handling of multiple requests for throughput",
        "• Sophisticated scheduling for background operations",
        "§Memory Optimization:",
        "• Automatic memory reclamation for resource efficiency",
        "• Proactive memory usage surveillance with alerts",
        "• Strategic invalidation of cached data for freshness",
        "• Careful resource tracking to prevent memory leaks",
        "§Reliability Engineering:",
        "• Heartbeat system to detect connection issues immediately",
        "• Self-healing reconnection with intelligent backoff strategy",
        "• Detailed error capturing with context for troubleshooting",
        "• Clean shutdown procedures to prevent data loss",
        "§Background Task Management:",
        "• Periodic optimization of metrics storage for performance",
        "• Continuous surveillance of system resource utilization",
        "• Regular checks of all external service connections",
        "• Self-managed recovery procedures for system health"
    ]
    create_content_slide("Technical Excellence", "Built for Performance and Reliability", technical_content)
    
    # Getting Started Slide
    getting_started_content = [
        "§Installation:",
        "• Single-click Discord authorization process",
        "• Default settings work immediately out of the box",
        "• No technical knowledge required for basic operation",
        "§Basic Configuration:",
        "• Use `/settings status` to view current configuration",
        "• Switch operating modes with `/settings mode` command",
        "• Add specific tracking channels with `/settings add-channel`",
        "§Testing the Bot:",
        "• Post any token contract address in configured channels",
        "• Try common ticker symbols like $SOL, $ETH, or $BTC",
        "• Use `/health` command (admin only) to check system status",
        "§Advanced Usage:",
        "• Configure different behavior for trading vs. discussion channels",
        "• Use `/github-checker` to analyze project repositories",
        "• Track first-caller performance to encourage research sharing"
    ]
    create_content_slide("Getting Started", "Simple Setup for Immediate Value", getting_started_content)
    
    # Contact & Support Slide
    contact_content = [
        "§Official Resources:",
        "• Website: [Your Website URL]",
        "• GitHub Repository: [Your GitHub Repo URL]",
        "• Documentation: [Your Docs URL]",
        "§Community Support:",
        "• Discord Server: [Your Support Server Invite]",
        "• Feature Requests: [Feature Request Form]",
        "• Bug Reports: [GitHub Issues Page]",
        "§Developer Contact:",
        "• Email: [Contact Email]",
        "• Twitter: [@YourHandle]"
    ]
    create_content_slide("Contact & Support", "Resources for Additional Help", contact_content)
    
    # Conclusion slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Apply background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['discord_blue']
    
    # Add title and subtitle
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Thank You!"
    subtitle.text = "TOPG Bot - The Ultimate Crypto Companion for Discord"
    
    # Format title and subtitle
    title_text_frame = title.text_frame
    title_paragraph = title_text_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_run = title_paragraph.runs[0]
    title_run.font.size = title_font_size
    title_run.font.color.rgb = colors['text_primary']
    title_run.font.bold = True
    
    subtitle_text_frame = subtitle.text_frame
    subtitle_paragraph = subtitle_text_frame.paragraphs[0]
    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    subtitle_run = subtitle_paragraph.runs[0]
    subtitle_run.font.size = subtitle_font_size
    subtitle_run.font.color.rgb = colors['text_primary']
    
    # Save the presentation
    prs.save(output_path)
    print(f"Presentation successfully created at: {os.path.abspath(output_path)}")
    return prs

if __name__ == "__main__":
    create_topg_presentation()